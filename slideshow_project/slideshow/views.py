# Create your views here.

from django.shortcuts import render, redirect
from .models import Image

def home(request):
    return redirect('slideshow') 
def slideshow(request):
    images = Image.objects.all()
    return render(request, 'slideshow/slideshow.html', {'images': images})

def image_detail(request, image_id):
    image = Image.objects.get(id=image_id)
    return render(request, 'slideshow/image_detail.html', {'image': image})

# slideshow/views.py


import os
from django.conf import settings
from . import utils
from django.shortcuts import render
from .forms import PresentationUploadForm
from .utils import parse_presentation, save_slide_images

from pptx import Presentation


def upload_presentation(request):
    if request.method == 'POST':
        form = PresentationUploadForm(request.POST, request.FILES)
        if form.is_valid():
            presentation_file = request.FILES['presentation']
            
            # Create the directory if it doesn't exist
            presentation_directory = os.path.join('media', 'presentations')
            os.makedirs(presentation_directory, exist_ok=True)
            
            # Save the uploaded file
            file_path = os.path.join(presentation_directory, presentation_file.name)
            with open(file_path, 'wb') as destination:
                for chunk in presentation_file.chunks():
                    destination.write(chunk)
            
            # Parse the uploaded presentation file
            presentation = Presentation(file_path)
            slide_images = []
            for idx, slide in enumerate(presentation.slides):
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # 13 is the code for an image shape
                        image = shape.image
                        image_bytes = image.blob
                        image_path = os.path.join('media', 'slides', f'slide_{idx}.png')
                        with open(image_path, 'wb') as img_file:
                            img_file.write(image_bytes)
                        slide_images.append(image_path)
            
            # Render the slideshow template with slide images
            return render(request, 'slideshow/slideshow.html', {'slide_images': slide_images})
    else:
        form = PresentationUploadForm()
    return render(request, 'slideshow/upload_presentation.html', {'form': form})